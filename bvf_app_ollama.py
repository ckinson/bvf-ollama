# bvf_app_ollama_utility_sectorized.py
# Streamlit BVF Builder (Sector-Smart Utility Layout) using local Ollama or OpenAI API
#
# New:
# - PPTX slide size picker: 16:9, 16:10, 4:3, A4 Landscape
# - Theme color pickers (affect visual, PDF, and PPTX)
# - Keeps: OpenAI v1/v0 compatibility, Ollama support, PDF export, editable PPTX export
#
# Extra dependency:
#   pip install python-pptx
#
# Other deps:
#   pip install streamlit ollama openai python-dotenv requests beautifulsoup4 lxml readability-lxml lxml-html-clean pdfminer.six plotly pandas pillow python-docx reportlab kaleido

import io
import json
from dataclasses import dataclass, field, asdict
from typing import List, Dict, Optional

import streamlit as st
import requests
from bs4 import BeautifulSoup
from readability import Document as ReadabilityDocument
from pdfminer.high_level import extract_text as pdf_extract_text
import pandas as pd
import plotly.graph_objects as go
from ollama import chat as ollama_chat
from docx import Document as DocxDocument
from reportlab.platypus import SimpleDocTemplate, Image as RLImage
from reportlab.lib.pagesizes import A4, landscape, portrait
from reportlab.lib.styles import getSampleStyleSheet
from PIL import Image as PILImage

# ---------------------------
# Streamlit & Colors
# ---------------------------
st.set_page_config(page_title="BVF Builder (Sector Smart • Ollama/OpenAI)", layout="wide")

PALETTE_DEFAULT = {
    "bg": "#FFFFFF",
    "exec_band": "#F2F2F2",
    "fin_band": "#F7F7F7",
    "functions_band_label": "#1F4E79",
    "function_tile": "#2E75B6",
    "function_body": "#EAF2FB",
    "kpi_band": "#DDEBF7",
    "priorities_band_label": "#833C99",
    "priority_tile": "#9E57B3",
    "priority_body": "#F3E9F8",
    "text_dark": "#0F172A",
}

def palette_ui(default: Dict[str, str]) -> Dict[str, str]:
    """Render color pickers and return a palette dict."""
    pal = dict(default)
    with st.expander("🎨 Theme colors (optional)", expanded=False):
        use_custom = st.checkbox("Customize theme colors", value=False)
        if use_custom:
            c1, c2, c3 = st.columns(3)
            with c1:
                pal["exec_band"] = st.color_picker("Executive band", pal["exec_band"])
                pal["functions_band_label"] = st.color_picker("Functions label band", pal["functions_band_label"])
                pal["function_tile"] = st.color_picker("Function header tile", pal["function_tile"])
                pal["kpi_band"] = st.color_picker("Operating KPIs band", pal["kpi_band"])
            with c2:
                pal["fin_band"] = st.color_picker("Financial/Operational band", pal["fin_band"])
                pal["function_body"] = st.color_picker("Function body", pal["function_body"])
                pal["priorities_band_label"] = st.color_picker("Priorities label band", pal["priorities_band_label"])
                pal["priority_tile"] = st.color_picker("Priority header tile", pal["priority_tile"])
            with c3:
                pal["priority_body"] = st.color_picker("Priority body", pal["priority_body"])
                pal["bg"] = st.color_picker("Background", pal["bg"])
                pal["text_dark"] = st.color_picker("Text color", pal["text_dark"])
    return pal

# ---------------------------
# Sector labels mapping
# ---------------------------
SECTORS = [
    "Auto-detect from content",
    "Utilities / Energy",
    "Retail",
    "Insurance",
    "Banking",
    "Telecom",
    "Manufacturing",
    "Healthcare",
    "Public Sector",
    "Technology / SaaS",
]

def get_sector_labels(sector: str) -> Dict[str, str]:
    dflt = {
        "exec_label": "Executive KPIs",
        "fin_label": "Financial / Operational KPIs",
        "functions_label": "Business Processes & Functions",
        "op_kpis_label": "Operating KPIs",
        "priorities_label": "Business Priorities",
        "tech_priorities_label": "Technology Priorities",
    }
    mapping = {
        "Utilities / Energy": dflt,
        "Retail": {
            "exec_label": "Executive KPIs",
            "fin_label": "Financial / Operational KPIs",
            "functions_label": "Value Chain & Functions",
            "op_kpis_label": "Store / Channel KPIs",
            "priorities_label": "Growth & Customer Priorities",
            "tech_priorities_label": "Digital & Technology Enablers",
        },
        "Insurance": {
            "exec_label": "Executive KPIs",
            "fin_label": "Financial / Operational KPIs",
            "functions_label": "Value Chain & Functions",
            "op_kpis_label": "Operational KPIs",
            "priorities_label": "Strategic Priorities",
            "tech_priorities_label": "Technology Enablers",
        },
        "Banking": {
            "exec_label": "Executive KPIs",
            "fin_label": "Financial / Operational KPIs",
            "functions_label": "Value Chain & Functions",
            "op_kpis_label": "Operational KPIs",
            "priorities_label": "Transformation Priorities",
            "tech_priorities_label": "Technology Enablers",
        },
        "Telecom": {
            "exec_label": "Executive KPIs",
            "fin_label": "Financial / Operational KPIs",
            "functions_label": "Network & Customer Functions",
            "op_kpis_label": "Operational KPIs",
            "priorities_label": "Growth & Network Priorities",
            "tech_priorities_label": "Technology Enablers",
        },
        "Manufacturing": {
            "exec_label": "Executive KPIs",
            "fin_label": "Financial / Operational KPIs",
            "functions_label": "Operations & Supply Chain Functions",
            "op_kpis_label": "Operational KPIs",
            "priorities_label": "Manufacturing Priorities",
            "tech_priorities_label": "Industry 4.0 Enablers",
        },
        "Healthcare": {
            "exec_label": "Executive KPIs",
            "fin_label": "Financial / Operational KPIs",
            "functions_label": "Clinical & Operational Functions",
            "op_kpis_label": "Clinical / Operational KPIs",
            "priorities_label": "Clinical & Transformation Priorities",
            "tech_priorities_label": "Digital Health Enablers",
        },
        "Public Sector": {
            "exec_label": "Mission Outcomes / KPIs",
            "fin_label": "Financial / Operational KPIs",
            "functions_label": "Programs & Service Functions",
            "op_kpis_label": "Performance KPIs",
            "priorities_label": "Policy & Service Priorities",
            "tech_priorities_label": "GovTech Enablers",
        },
        "Technology / SaaS": {
            "exec_label": "Executive KPIs",
            "fin_label": "Financial / Product KPIs",
            "functions_label": "Product & GTM Functions",
            "op_kpis_label": "Product / Ops KPIs",
            "priorities_label": "Growth & Platform Priorities",
            "tech_priorities_label": "Platform & Engineering Enablers",
        },
    }
    return mapping.get(sector, dflt)

# ---------------------------
# Data Model
# ---------------------------
@dataclass
class BVF:
    company: str
    executive_kpis: List[str] = field(default_factory=list)
    financial_operational_kpis: List[str] = field(default_factory=list)
    business_functions: List[str] = field(default_factory=list)
    operating_kpis_by_function: Dict[str, List[str]] = field(default_factory=dict)
    function_projects: Dict[str, List[str]] = field(default_factory=dict)
    business_priorities: List[str] = field(default_factory=list)
    technology_priorities_by_business_priority: Dict[str, List[str]] = field(default_factory=dict)
    sources: List[str] = field(default_factory=list)

    def curate(self):
        self.executive_kpis = sorted(set([s for s in self.executive_kpis if s.strip()]))[:8]
        self.financial_operational_kpis = sorted(set([s for s in self.financial_operational_kpis if s.strip()]))[:10]
        seen = set(); ordered = []
        for f in self.business_functions:
            if f and f not in seen:
                ordered.append(f); seen.add(f)
        self.business_functions = ordered[:10] if ordered else list(self.operating_kpis_by_function.keys())[:10]
        self.operating_kpis_by_function = {k: sorted(set([x for x in v if x.strip()]))[:8] for k, v in self.operating_kpis_by_function.items()}
        self.function_projects = {k: sorted(set([x for x in v if x.strip()]))[:6] for k, v in self.function_projects.items()}
        self.business_priorities = [*dict.fromkeys([s for s in self.business_priorities if s.strip()])][:8]
        self.technology_priorities_by_business_priority = {k: sorted(set([x for x in v if x.strip()]))[:8] for k, v in self.technology_priorities_by_business_priority.items()}

    def to_frame(self) -> pd.DataFrame:
        rows = []
        for f in self.business_functions:
            rows.append({"Section": "Function Projects", "Lane": f, "Items": "\n".join(self.function_projects.get(f, []))})
        for f, items in self.operating_kpis_by_function.items():
            rows.append({"Section": "Operating KPIs", "Lane": f, "Items": "\n".join(items)})
        for p in self.business_priorities:
            rows.append({"Section": "Business Priority", "Lane": p, "Items": "\n".join(self.technology_priorities_by_business_priority.get(p, []))})
        return pd.DataFrame(rows)

# ---------------------------
# Helpers: ingest
# ---------------------------
def fetch_text(url: str) -> str:
    try:
        if url.lower().endswith(".pdf"):
            return pdf_extract_text(url)
        resp = requests.get(url, timeout=30, headers={"User-Agent": "Mozilla/5.0"})
        resp.raise_for_status()
        if "application/pdf" in resp.headers.get("Content-Type", ""):
            with open("temp.pdf", "wb") as fh:
                fh.write(resp.content)
            return pdf_extract_text("temp.pdf")
        doc = ReadabilityDocument(resp.text)
        soup = BeautifulSoup(doc.summary(), "lxml")
        return soup.get_text(separator=" ", strip=True)
    except Exception:
        return ""

def read_uploaded_file(uploaded_file) -> str:
    name = uploaded_file.name.lower()
    if name.endswith(".pdf"):
        with open("temp_uploaded.pdf", "wb") as f:
            f.write(uploaded_file.read())
        return pdf_extract_text("temp_uploaded.pdf")
    elif name.endswith(".docx"):
        doc = DocxDocument(uploaded_file)
        return "\n".join([p.text for p in doc.paragraphs if p.text.strip()])
    else:
        return uploaded_file.read().decode("utf-8", errors="ignore")

# ---------------------------
# LLM: Providers (Ollama/OpenAI) + JSON parsing
# ---------------------------
def call_ollama(messages: List[Dict], model: str) -> str:
    resp = ollama_chat(model=model, messages=messages)
    return resp["message"]["content"]

# ---- OpenAI compatibility (v1.x AND legacy v0.x) ----
def _openai_chat_v1(messages: List[Dict], model: str, api_key: str, **kwargs) -> str:
    from openai import OpenAI  # v1.x
    client = OpenAI(api_key=api_key)
    out = client.chat.completions.create(
        model=model,
        messages=messages,
        temperature=kwargs.get("temperature", 0.2),
    )
    return out.choices[0].message.content

def _openai_chat_v0(messages: List[Dict], model: str, api_key: str, **kwargs) -> str:
    import openai  # legacy v0.x
    openai.api_key = api_key
    out = openai.ChatCompletion.create(  # type: ignore[attr-defined]
        model=model,
        messages=messages,
        temperature=kwargs.get("temperature", 0.2),
    )
    return out["choices"][0]["message"]["content"]

def call_openai_compat(messages: List[Dict], model: str, api_key: str, **kwargs) -> str:
    try:
        return _openai_chat_v1(messages, model, api_key, **kwargs)
    except ImportError:
        return _openai_chat_v0(messages, model, api_key, **kwargs)
    except Exception as e:
        st.error(f"OpenAI error: {e}")
        return ""

def llm_generate_json_text(prompt: str, provider: str, model: str, api_key: Optional[str]) -> str:
    system = {"role": "system", "content": "You are a senior business strategist. Return only valid JSON with no extra text."}
    user = {"role": "user", "content": prompt}
    if provider == "OpenAI API":
        if not api_key:
            st.error("Please enter your OpenAI API key.")
            return ""
        return call_openai_compat([system, user], model=model, api_key=api_key, temperature=0.2)
    else:
        return call_ollama([system, user], model=model)

def parse_json_safely(output: str) -> Optional[dict]:
    try:
        return json.loads(output)
    except json.JSONDecodeError:
        start = output.find("{"); end = output.rfind("}")
        if start != -1 and end != -1:
            try:
                return json.loads(output[start:end+1])
            except json.JSONDecodeError:
                return None
        return None

def detect_sector(corp_text: str, provider: str, model: str, api_key: Optional[str]) -> str:
    options = [s for s in SECTORS if s != "Auto-detect from content"]
    prompt = f"""You are an industry classifier. Choose the single best-fitting sector label from this list:
{options}
Return ONLY the label string, nothing else.

Text:
{corp_text[:6000]}
"""
    system = {"role": "system", "content": "Return only one label string, nothing else."}
    user = {"role": "user", "content": prompt}
    if provider == "OpenAI API":
        if not api_key:
            st.error("Please enter your OpenAI API key.")
            return "Utilities / Energy"
        try:
            choice = call_openai_compat([system, user], model=model, api_key=api_key, temperature=0.0).strip()
        except Exception as e:
            st.error(f"OpenAI error (sector detect): {e}")
            return "Utilities / Energy"
    else:
        choice = call_ollama([system, user], model=model).strip()
    return choice if choice in options else "Utilities / Energy"

# ---------------------------
# Extraction (sector-agnostic schema)
# ---------------------------
def extract_bvf(company: str, corp_text: str, provider: str, model: str, api_key: Optional[str]) -> BVF:
    schema_hint = json.dumps({
        "executive_kpis": ["string"],
        "financial_operational_kpis": ["string"],
        "business_functions": ["string"],
        "operating_kpis_by_function": {"FunctionName": ["KPI1", "KPI2"]},
        "function_projects": {"FunctionName": ["Project1", "Project2"]},
        "business_priorities": ["string"],
        "technology_priorities_by_business_priority": {"BusinessPriorityName": ["Tech priority 1", "Tech priority 2"]},
        "sources": ["string"]
    }, indent=2)

    prompt = f"""Company: {company}

GOAL
Produce a curated Business Value Framework (BVF) using the layers below.
Use concise, sector-specific, non-generic language. Remove duplicates and prioritize impact.

LAYOUT & FIELDS (STRICT JSON ONLY):
{schema_hint}

GUIDANCE
- business_functions should reflect end-to-end operations for the client's sector.
- operating_kpis_by_function: include 3–8 crisp, quantifiable KPIs per function.
- function_projects (optional): 3–6 short bullets reflecting realistic initiatives.
- business_priorities: 5–8 transformation themes.
- technology_priorities_by_business_priority: 4–8 technology levers per priority.

RULES
- Return ONLY valid JSON. No explanations, code fences, or extra text.

TEXT TO ANALYZE
{corp_text[:100000]}
"""

    output = llm_generate_json_text(prompt, provider=provider, model=model, api_key=api_key)
    if not output:
        return BVF(company=company)

    data = parse_json_safely(output)
    if not data:
        st.error("❌ Model did not return valid JSON after salvage.")
        return BVF(company=company)

    bvf = BVF(company=company, **data)
    bvf.curate()
    return bvf

# ---------------------------
# Visualization helpers (rounded rectangles + layout)
# ---------------------------
def bulletify(items: List[str]) -> str:
    if not items:
        return "—"
    return "<br>".join([f"• {x}" for x in items])

def _rounded_rect_path(x0, y0, x1, y1, r):
    r = max(0.0, min(r, (x1 - x0) / 2.0, (y1 - y0) / 2.0))
    return (
        f"M {x0+r},{y0} "
        f"L {x1-r},{y0} "
        f"Q {x1},{y0} {x1},{y0+r} "
        f"L {x1},{y1-r} "
        f"Q {x1},{y1} {x1-r},{y1} "
        f"L {x0+r},{y1} "
        f"Q {x0},{y1} {x0},{y1-r} "
        f"L {x0},{y0+r} "
        f"Q {x0},{y0} {x0+r},{y0} Z"
    )

def add_roundrect(fig, x0, y0, x1, y1, radius, fill, line, width=1):
    fig.add_shape(
        type="path",
        path=_rounded_rect_path(x0, y0, x1, y1, radius),
        line=dict(color=line, width=width),
        fillcolor=fill,
        layer="below",
    )

def text_center(fig, x, y, html, size=14, color="#000000"):
    fig.add_annotation(x=x, y=y, text=html, showarrow=False, yanchor="middle",
                       font=dict(size=size, color=color))

def render_bvf_figure_utility_layout(bvf: BVF, sector_labels: Dict[str, str], palette: Dict[str, str], height_px: int = 1300) -> go.Figure:
    functions = bvf.business_functions or list(bvf.operating_kpis_by_function.keys())
    functions = functions[:10] if functions else ["Function"]
    n_cols = max(6, len(functions))

    # Taller rows + small vertical gaps between layers
    ROW_EXEC = 1.4
    ROW_FIN = 1.4
    ROW_LABEL = 0.7
    ROW_FUNCTIONS = 2.9
    ROW_OP_KPIS = 2.9
    ROW_PRIORITIES_LABEL = 0.7
    ROW_PRIORITIES = 3.2
    GAP = 0.18  # vertical gap between layers
    TITLE_H = 0.9

    content_rows = (
        ROW_EXEC + GAP +
        ROW_FIN + GAP +
        ROW_LABEL +
        ROW_FUNCTIONS + GAP +
        ROW_OP_KPIS + GAP +
        ROW_PRIORITIES_LABEL +
        ROW_PRIORITIES
    )
    total_with_title = TITLE_H + content_rows

    fig = go.Figure()
    fig.update_xaxes(visible=False, range=[0, n_cols])
    fig.update_yaxes(visible=False, range=[0, total_with_title])
    fig.update_layout(
        height=height_px,
        margin=dict(l=30, r=30, t=50, b=30),
        plot_bgcolor=palette["bg"],
        paper_bgcolor=palette["bg"],
        showlegend=False,
    )

    R_SMALL = 0.10
    R_MED = 0.16
    R_LARGE = 0.22

    y = total_with_title
    # Title
    y -= TITLE_H
    text_center(fig, n_cols/2, y + TITLE_H/2, "<b>Business Value Framework</b>", size=22, color=palette["text_dark"])

    # Exec KPIs
    y -= ROW_EXEC
    add_roundrect(fig, 0, y, n_cols, y+ROW_EXEC, R_LARGE, palette["exec_band"], palette["exec_band"], width=1)
    text_center(fig, n_cols/2, y + ROW_EXEC/2, f"<b>{sector_labels['exec_label']}</b><br><br>{bulletify(bvf.executive_kpis)}", color=palette["text_dark"])
    y -= GAP

    # Fin/Op KPIs
    y -= ROW_FIN
    add_roundrect(fig, 0, y, n_cols, y+ROW_FIN, R_LARGE, palette["fin_band"], palette["fin_band"], width=1)
    text_center(fig, n_cols/2, y + ROW_FIN/2, f"<b>{sector_labels['fin_label']}</b><br><br>{bulletify(bvf.financial_operational_kpis)}", color=palette["text_dark"])
    y -= GAP

    # Functions label
    y -= ROW_LABEL
    add_roundrect(fig, 0, y, n_cols, y+ROW_LABEL, R_SMALL, palette["functions_band_label"], palette["functions_band_label"], width=0)
    text_center(fig, 0.7, y + ROW_LABEL/2, f"<b style='color:white'>{sector_labels['functions_label']}</b>", color="white")

    # Function tiles
    y -= ROW_FUNCTIONS
    for i, f in enumerate(functions):
        x0 = i*(n_cols/len(functions)); x1 = (i+1)*(n_cols/len(functions))
        add_roundrect(fig, x0, y+ROW_FUNCTIONS*0.78, x1, y+ROW_FUNCTIONS, R_MED, palette["function_tile"], palette["function_tile"], width=0)
        text_center(fig, (x0+x1)/2, y+ROW_FUNCTIONS*0.89, f"<b style='color:white'>{f}</b>", size=13, color="white")
        add_roundrect(fig, x0, y, x1, y+ROW_FUNCTIONS*0.78, R_MED, palette["function_body"], palette["function_body"], width=1)
        bullets = bvf.function_projects.get(f, [])
        text_center(fig, (x0+x1)/2, y+ROW_FUNCTIONS*0.39, bulletify(bullets), size=12, color=palette["text_dark"])
    y -= GAP

    # Operating KPIs per function
    y -= ROW_OP_KPIS
    add_roundrect(fig, 0, y, n_cols, y+ROW_OP_KPIS, R_SMALL, palette["kpi_band"], palette["kpi_band"], width=1)
    for i, f in enumerate(functions):
        x0 = i*(n_cols/len(functions)); x1 = (i+1)*(n_cols/len(functions))
        add_roundrect(fig, x0+0.06, y+0.06, x1-0.06, y+ROW_OP_KPIS-0.06, R_SMALL, "#FFFFFF", "#CBD5E1", width=1)
        kp = bvf.operating_kpis_by_function.get(f, [])
        text_center(fig, (x0+x1)/2, y+ROW_OP_KPIS/2, f"<b>{f} — {sector_labels['op_kpis_label']}</b><br><br>{bulletify(kp)}", size=12, color=palette["text_dark"])
    y -= GAP

    # Priorities label
    y -= ROW_PRIORITIES_LABEL
    add_roundrect(fig, 0, y, n_cols, y+ROW_PRIORITIES_LABEL, R_SMALL, palette["priorities_band_label"], palette["priorities_band_label"], width=0)
    text_center(fig, 0.8, y + ROW_PRIORITIES_LABEL/2, f"<b style='color:white'>{sector_labels['priorities_label']}</b>", color="white")

    # Priority tiles
    y -= ROW_PRIORITIES
    priorities = bvf.business_priorities or list(bvf.technology_priorities_by_business_priority.keys())
    if not priorities:
        priorities = ["Priority"]
    if len(priorities) > len(functions):
        priorities = priorities[:len(functions)]
    for i, p in enumerate(priorities):
        x0 = i*(n_cols/len(priorities)); x1 = (i+1)*(n_cols/len(priorities))
        add_roundrect(fig, x0, y, x1, y+ROW_PRIORITIES, R_MED, palette["priority_body"], palette["priority_body"], width=1)
        add_roundrect(fig, x0, y+ROW_PRIORITIES*0.78, x1, y+ROW_PRIORITIES, R_MED, palette["priority_tile"], palette["priority_tile"], width=0)
        text_center(fig, (x0+x1)/2, y+ROW_PRIORITIES*0.89, f"<b style='color:white'>{p}</b>", size=13, color="white")
        techs = bvf.technology_priorities_by_business_priority.get(p, [])
        body = f"<b>{sector_labels['tech_priorities_label']}</b><br><br>{bulletify(techs)}"
        text_center(fig, (x0+x1)/2, y+ROW_PRIORITIES*0.39, body, size=12, color=palette["text_dark"])

    return fig

# ---------------------------
# PDF Export (visual table)
# ---------------------------
def export_visual_pdf(fig: go.Figure, filename: str, orientation: str = "Landscape"):
    try:
        png_bytes = fig.to_image(format="png", width=2400, height=1400, scale=2)  # crisp export
    except Exception:
        st.error("Image export failed. Ensure 'kaleido' is installed: pip install kaleido")
        raise

    page_size = landscape(A4) if orientation.lower().startswith("land") else portrait(A4)
    page_w, page_h = page_size
    margin = 24  # points

    img = PILImage.open(io.BytesIO(png_bytes))
    iw, ih = img.size
    max_w = page_w - 2 * margin
    max_h = page_h - 2 * margin
    scale = min(max_w / iw, max_h / ih)
    draw_w = iw * scale
    draw_h = ih * scale

    styles = getSampleStyleSheet()
    doc = SimpleDocTemplate(filename, pagesize=page_size,
                            leftMargin=margin, rightMargin=margin,
                            topMargin=margin, bottomMargin=margin)
    story = []
    rl_img = RLImage(io.BytesIO(png_bytes), width=draw_w, height=draw_h)
    story.append(rl_img)
    doc.build(story)

# ---------------------------
# NEW: PPTX Export (fully editable) with slide ratio + theme colors
# ---------------------------
def export_visual_pptx(bvf: BVF, sector_labels: Dict[str, str], filename: str, palette: Dict[str, str], slide_ratio: str):
    try:
        from pptx import Presentation
        from pptx.util import Inches, Pt
        from pptx.enum.text import PP_ALIGN
        from pptx.enum.text import MSO_ANCHOR
        from pptx.dml.color import RGBColor
        from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
    except Exception as e:
        st.error(f"PPTX export requires python-pptx. Install it with: pip install python-pptx\nError: {e}")
        raise

    def hex_to_rgb(hex_color: str):
        hex_color = hex_color.lstrip("#")
        return tuple(int(hex_color[i:i+2], 16) for i in (0, 2, 4))

    # Slide size presets (in inches)
    if "16:9" in slide_ratio:
        slide_w_in, slide_h_in = 13.33, 7.5
    elif "16:10" in slide_ratio:
        slide_w_in, slide_h_in = 12.8, 8.0
    elif "4:3" in slide_ratio:
        slide_w_in, slide_h_in = 10.0, 7.5
    else:  # A4 Landscape
        slide_w_in, slide_h_in = 11.69, 8.27

    prs = Presentation()
    prs.slide_width = Inches(slide_w_in)
    prs.slide_height = Inches(slide_h_in)
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank

    # Layout metrics (inches)
    MARGIN_IN = 0.35
    EMU_PER_INCH = 914400
    slide_width_in = prs.slide_width / EMU_PER_INCH
    slide_height_in = prs.slide_height / EMU_PER_INCH
    width_in = slide_width_in - 2 * MARGIN_IN
    height_in = slide_height_in - 2 * MARGIN_IN

    # Row proportions matching chart layout
    ROW_EXEC = 1.4
    ROW_FIN = 1.4
    ROW_LABEL = 0.7
    ROW_FUNCTIONS = 2.9
    ROW_OP_KPIS = 2.9
    ROW_PRIORITIES_LABEL = 0.7
    ROW_PRIORITIES = 3.2
    GAP = 0.18
    TITLE_H = 0.9

    unit_total = TITLE_H + (ROW_EXEC + GAP + ROW_FIN + GAP + ROW_LABEL + ROW_FUNCTIONS + GAP + ROW_OP_KPIS + GAP + ROW_PRIORITIES_LABEL + ROW_PRIORITIES)
    unit_to_in = height_in / unit_total  # inches per "unit"

    # Helper to add rounded rectangle
    def add_box(xu, yu, wu, hu, fill_hex, line_hex=None, corner=0.2):
        left = Inches(MARGIN_IN + xu * unit_to_in)
        top = Inches(MARGIN_IN + yu * unit_to_in)
        w = Inches(wu * unit_to_in)
        h = Inches(hu * unit_to_in)
        shape = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, left, top, w, h)
        if line_hex is None:
            line_hex = fill_hex
        r, g, b = hex_to_rgb(fill_hex); shape.fill.solid(); shape.fill.fore_color.rgb = RGBColor(r, g, b)
        lr, lg, lb = hex_to_rgb(line_hex); shape.line.color.rgb = RGBColor(lr, lg, lb); shape.line.width = Pt(1.25)
        try:
            shape.adjustments[0] = corner
        except Exception:
            pass
        return shape

    def set_text(shape, lines: List[str], bold_first=False, center=True, font_size=12, font_color=palette["text_dark"]):
        tf = shape.text_frame
        tf.clear()
        fr, fg, fb = hex_to_rgb(font_color)
        # first line
        p = tf.paragraphs[0]
        run = p.add_run()
        run.text = lines[0] if lines else ""
        run.font.size = Pt(font_size)
        run.font.bold = bold_first
        run.font.color.rgb = RGBColor(fr, fg, fb)
        p.space_after = Pt(4)
        if center:
            p.alignment = PP_ALIGN.CENTER
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        # rest
        for ln in lines[1:]:
            p2 = tf.add_paragraph()
            p2.text = ln
            p2.level = 0
            p2.space_after = Pt(2)
            p2.font.size = Pt(font_size)
            p2.font.color.rgb = RGBColor(fr, fg, fb)
            if center:
                p2.alignment = PP_ALIGN.CENTER

    def bulletize(items: List[str]) -> List[str]:
        if not items:
            return ["—"]
        return [f"• {x}" for x in items]

    # Compute columns (functions & priorities)
    functions = bvf.business_functions or list(bvf.operating_kpis_by_function.keys())
    functions = functions[:10] if functions else ["Function"]
    fn_cols = len(functions)
    priorities = bvf.business_priorities or list(bvf.technology_priorities_by_business_priority.keys())
    if not priorities:
        priorities = ["Priority"]
    if len(priorities) > fn_cols:
        priorities = priorities[:fn_cols]
    pr_cols = len(priorities)

    total_units_w = width_in / unit_to_in
    col_w = total_units_w / fn_cols
    col_w_p = total_units_w / pr_cols

    # Current y position in units from top margin
    y = 0.0

    # Title (text only)
    title_shape = add_box(0, y, total_units_w, TITLE_H, palette["bg"], palette["bg"], corner=0.0)
    set_text(title_shape, ["Business Value Framework"], bold_first=True, center=True, font_size=22)
    y += TITLE_H

    # Exec KPIs band
    exec_shape = add_box(0, y, total_units_w, ROW_EXEC, palette["exec_band"], palette["exec_band"])
    exec_lines = [f"{sector_labels['exec_label']}"] + [""] + bulletize(bvf.executive_kpis)
    set_text(exec_shape, exec_lines, bold_first=True, center=True, font_size=12)
    y += ROW_EXEC + GAP

    # Fin/Op KPIs band
    fin_shape = add_box(0, y, total_units_w, ROW_FIN, palette["fin_band"], palette["fin_band"])
    fin_lines = [f"{sector_labels['fin_label']}"] + [""] + bulletize(bvf.financial_operational_kpis)
    set_text(fin_shape, fin_lines, bold_first=True, center=True, font_size=12)
    y += ROW_FIN + GAP

    # Functions label band
    flabel_shape = add_box(0, y, total_units_w, ROW_LABEL, palette["functions_band_label"], palette["functions_band_label"])
    set_text(flabel_shape, [f"{sector_labels['functions_label']}"], bold_first=True, center=True, font_size=12, font_color="#FFFFFF")
    y += ROW_LABEL

    # Function tiles (header strip + body)
    f_body_h = ROW_FUNCTIONS * 0.78
    f_header_h = ROW_FUNCTIONS - f_body_h
    for i, f in enumerate(functions):
        x = i * col_w
        h_shape = add_box(x, y + f_body_h, col_w, f_header_h, palette["function_tile"], palette["function_tile"], corner=0.15)
        set_text(h_shape, [f], bold_first=True, center=True, font_size=12, font_color="#FFFFFF")
        b_shape = add_box(x, y, col_w, f_body_h, palette["function_body"], palette["function_body"], corner=0.15)
        set_text(b_shape, bulletize(bvf.function_projects.get(f, [])), bold_first=False, center=True, font_size=11)
    y += ROW_FUNCTIONS + GAP

    # Operating KPIs band
    kpi_band = add_box(0, y, total_units_w, ROW_OP_KPIS, palette["kpi_band"], palette["kpi_band"], corner=0.10)
    inner_pad = 0.06
    for i, f in enumerate(functions):
        x = i * col_w
        card = add_box(x + inner_pad, y + inner_pad, col_w - 2*inner_pad, ROW_OP_KPIS - 2*inner_pad, "#FFFFFF", "#CBD5E1", corner=0.10)
        lines = [f"{f} — {sector_labels['op_kpis_label']}"] + [""] + bulletize(bvf.operating_kpis_by_function.get(f, []))
        set_text(card, lines, bold_first=True, center=True, font_size=11)
    y += ROW_OP_KPIS + GAP

    # Priorities label band
    plabel_shape = add_box(0, y, total_units_w, ROW_PRIORITIES_LABEL, palette["priorities_band_label"], palette["priorities_band_label"])
    set_text(plabel_shape, [f"{sector_labels['priorities_label']}"], bold_first=True, center=True, font_size=12, font_color="#FFFFFF")
    y += ROW_PRIORITIES_LABEL

    # Priorities tiles
    p_body_h = ROW_PRIORITIES * 0.78
    p_header_h = ROW_PRIORITIES - p_body_h
    for i, p in enumerate(priorities):
        x = i * col_w_p
        p_body = add_box(x, y, col_w_p, p_body_h, palette["priority_body"], palette["priority_body"], corner=0.15)
        lines = [f"{sector_labels['tech_priorities_label']}"] + [""] + bulletize(bvf.technology_priorities_by_business_priority.get(p, []))
        set_text(p_body, lines, bold_first=True, center=True, font_size=11)
        p_head = add_box(x, y + p_body_h, col_w_p, p_header_h, palette["priority_tile"], palette["priority_tile"], corner=0.15)
        set_text(p_head, [p], bold_first=True, center=True, font_size=12, font_color="#FFFFFF")

    prs.save(filename)

# ---------------------------
# UI
# ---------------------------
st.title("🧭 BVF Builder — Sector Smart (Ollama / OpenAI)")

company = st.text_input("Company name", placeholder="e.g., Aviva Insurance")

provider = st.selectbox("LLM provider", ["Ollama (local)", "OpenAI API"], index=0)
if provider == "Ollama (local)":
    model_name = st.selectbox("Model", ["llama3", "mistral", "gemma", "qwen"], index=0)
    openai_api_key = None
else:
    model_name = st.selectbox("Model", ["gpt-4o-mini", "gpt-4.1-mini", "gpt-4o", "gpt-3.5-turbo"], index=0)
    openai_api_key = st.text_input("OpenAI API key", type="password", placeholder="sk-...", help="Your key is kept in memory only for this session.")

pdf_orientation = st.selectbox("PDF orientation", ["Landscape", "Portrait"], index=0)
pptx_ratio = st.selectbox("PPTX slide size", ["16:9 (Widescreen)", "16:10", "4:3 (Standard)", "A4 Landscape"], index=0)

selected_sector = st.selectbox("Industry sector", SECTORS, index=0)

manual_urls = st.text_area("Optional: paste specific URLs (one per line)", height=100, placeholder="https://example.com/strategy.pdf\nhttps://investors.example.com/annual-report")
manual_text = st.text_area("Paste raw strategy text here", height=200)

uploaded_files = st.file_uploader("Upload local PDF or DOCX strategy files", type=["pdf", "docx"], accept_multiple_files=True)

# Theme palette UI
palette = palette_ui(PALETTE_DEFAULT)

colA, colB, colC, colD = st.columns(4)
with colA:
    if st.button("Load Uploaded Files"):
        file_text = ""
        for file in uploaded_files or []:
            file_text += "\n\nSOURCE: " + file.name + "\n" + read_uploaded_file(file)
        manual_text += ("\n" + file_text) if file_text else ""
        st.session_state["ingested_text"] = manual_text

with colB:
    if st.button("Fetch from URLs"):
        all_text = ""
        for u in manual_urls.splitlines():
            u = u.strip()
            if u:
                all_text += f"\n\nSOURCE: {u}\n{fetch_text(u)}"
        manual_text += ("\n" + all_text) if all_text else ""
        st.session_state["ingested_text"] = manual_text

with colC:
    if st.button("Auto-detect sector"):
        full_text = st.session_state.get("ingested_text", "") or manual_text
        if not full_text.strip():
            st.warning("Add some text or URLs/files first so I can detect the sector.")
        else:
            sector = detect_sector(
                full_text,
                provider="OpenAI API" if provider == "OpenAI API" else "Ollama (local)",
                model=model_name,
                api_key=openai_api_key,
            )
            st.success(f"Detected sector: {sector}")
            st.session_state["sector"] = sector

with colD:
    build_disabled = not company or not (manual_text.strip() or st.session_state.get("ingested_text", "").strip())
    if st.button("Build BVF", disabled=build_disabled):
        full_text = st.session_state.get("ingested_text", "") or manual_text
        with st.spinner("Generating BVF..."):
            bvf = extract_bvf(
                company,
                full_text,
                provider=("OpenAI API" if provider == "OpenAI API" else "Ollama (local)"),
                model=model_name,
                api_key=openai_api_key,
            )
            st.session_state["bvf"] = bvf
            if st.session_state.get("sector"):
                st.session_state["sector_locked"] = st.session_state["sector"]

# Results
bvf: Optional[BVF] = st.session_state.get("bvf")
effective_sector = (
    st.session_state.get("sector_locked")
    or (st.session_state.get("sector") if selected_sector == "Auto-detect from content" else selected_sector)
)
if selected_sector != "Auto-detect from content":
    effective_sector = selected_sector
labels = get_sector_labels(effective_sector or "Utilities / Energy")

if bvf and (bvf.executive_kpis or bvf.business_functions):
    # Visual only (no on-screen JSON)
    st.subheader("Visual")
    fig = render_bvf_figure_utility_layout(bvf, labels, palette)
    st.plotly_chart(fig, use_container_width=True)

    # Exports
    st.subheader("Export")
    df = bvf.to_frame()
    st.download_button("Download JSON", json.dumps(asdict(bvf), indent=2), file_name=f"{bvf.company}_BVF.json")
    st.download_button("Download CSV", df.to_csv(index=False).encode("utf-8"), file_name=f"{bvf.company}_BVF.csv", mime="text/csv")

    # PDF (visual layout)
    pdf_filename = f"{bvf.company}_BVF_visual.pdf"
    try:
        export_visual_pdf(fig, pdf_filename, orientation=pdf_orientation)
        with open(pdf_filename, "rb") as pdf_file:
            st.download_button("Download PDF (visual layout)", pdf_file, file_name=pdf_filename, mime="application/pdf")
    except Exception:
        pass

    # PPTX (editable) with selected ratio + palette
    pptx_filename = f"{bvf.company}_BVF_editable.pptx"
    try:
        export_visual_pptx(bvf, labels, pptx_filename, palette, pptx_ratio)
        with open(pptx_filename, "rb") as f:
            st.download_button("Download PPTX (editable)", f, file_name=pptx_filename, mime="application/vnd.openxmlformats-officedocument.presentationml.presentation")
    except Exception:
        pass
